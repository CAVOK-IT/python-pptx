# encoding: utf-8

"""
Step implementations for slide master-related features
"""

from __future__ import absolute_import

from behave import given, then

from pptx import Presentation
from pptx.parts.slidelayout import SlideLayout
from pptx.parts.slidemaster import _MasterShapeTree, _SlideLayouts
from pptx.shapes.shape import BaseShape

from .helpers import test_pptx


# given ===================================================

@given('a master shape collection containing two shapes')
def given_master_shape_collection_containing_two_shapes(context):
    prs = Presentation(test_pptx('mst-shapes'))
    context.master_shapes = prs.slide_master.shapes


@given('a slide master having two shapes')
def given_slide_master_having_two_shapes(context):
    prs = Presentation(test_pptx('mst-shapes'))
    context.slide_master = prs.slide_master


@given('a slide master having two slide layouts')
def given_slide_master_having_two_layouts(context):
    prs = Presentation(test_pptx('mst-slide-layouts'))
    context.slide_master = prs.slide_master


@given('a slide layout collection containing two layouts')
def given_slide_layout_collection_containing_two_layouts(context):
    prs = Presentation(test_pptx('mst-slide-layouts'))
    context.slide_layouts = prs.slide_master.slide_layouts


# then ====================================================

@then('I can access a master shape by index')
def then_can_access_master_shape_by_index(context):
    master_shapes = context.master_shapes
    for idx in range(2):
        master_shape = master_shapes[idx]
        assert isinstance(master_shape, BaseShape)


@then('I can access a slide layout by index')
def then_can_access_slide_layout_by_index(context):
    slide_layouts = context.slide_layouts
    for idx in range(2):
        slide_layout = slide_layouts[idx]
        assert isinstance(slide_layout, SlideLayout)


@then('I can access the shape collection of the slide master')
def then_can_access_shape_collection_of_slide_master(context):
    slide_master = context.slide_master
    master_shapes = slide_master.shapes
    msg = 'SlideMaster.shapes not instance of _MasterShapeTree'
    assert isinstance(master_shapes, _MasterShapeTree), msg


@then('I can access the slide layouts of the slide master')
def then_can_access_slide_layouts_of_slide_master(context):
    slide_master = context.slide_master
    slide_layouts = slide_master.slide_layouts
    msg = 'SlideMaster.slide_layouts not instance of _SlideLayouts'
    assert isinstance(slide_layouts, _SlideLayouts), msg


@then('I can iterate over the master shapes')
def then_can_iterate_over_the_master_shapes(context):
    master_shapes = context.master_shapes
    actual_count = 0
    for master_shape in master_shapes:
        actual_count += 1
        assert isinstance(master_shape, BaseShape)
    assert actual_count == 2


@then('I can iterate over the slide layouts')
def then_can_iterate_over_the_slide_layouts(context):
    slide_layouts = context.slide_layouts
    actual_count = 0
    for slide_layout in slide_layouts:
        actual_count += 1
        assert isinstance(slide_layout, SlideLayout)
    assert actual_count == 2


@then('the length of the master shape collection is 2')
def then_len_of_master_shape_collection_is_2(context):
    slide_master = context.slide_master
    master_shapes = slide_master.shapes
    assert len(master_shapes) == 2, (
        'expected len(master_shapes) of 2, got %s' % len(master_shapes)
    )


@then('the length of the slide layout collection is 2')
def then_len_of_slide_layout_collection_is_2(context):
    slide_master = context.slide_master
    slide_layouts = slide_master.slide_layouts
    assert len(slide_layouts) == 2, (
        'expected len(slide_layouts) of 2, got %s' % len(slide_layouts)
    )
